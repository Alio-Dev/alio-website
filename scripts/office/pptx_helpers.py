"""Shared helpers for the Alio PowerPoint template (Presentation.pptx)."""

from pptx.util import Cm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

from brand import (
    FONT_DISPLAY, FONT_BODY, FONT_MONO,
    PRIMARY_700, PRIMARY_950, ACCENT_500, ACCENT_600,
    NEUTRAL_950, NEUTRAL_700, NEUTRAL_500, NEUTRAL_200, NEUTRAL_50, WHITE,
    hx,
)

RGB = lambda h: RGBColor(*hx(h))

SLIDE_W = Cm(33.867)  # 16:9 widescreen
SLIDE_H = Cm(19.05)
MARGIN = Cm(1.5)


def blank_layout(prs):
    return prs.slide_layouts[6]  # "Blank" in the default python-pptx template


def add_slide(prs):
    return prs.slides.add_slide(blank_layout(prs))


def set_bg(slide, hex_color):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGB(hex_color)


def add_gradient_rect(slide, left, top, width, height, angle=0):
    """Brand gradient (BRAND.md sec.3: linear-gradient(90deg,#07B7D1,#2B3990))."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    shape.shadow.inherit = False
    fill = shape.fill
    fill.gradient()
    stops = fill.gradient_stops
    stops[0].color.rgb = RGB(ACCENT_500)
    stops[0].position = 0.0
    stops[1].color.rgb = RGB(PRIMARY_700)
    stops[1].position = 1.0
    try:
        fill.gradient_angle = angle
    except Exception:
        pass
    return shape


def add_rect(slide, left, top, width, height, color, line=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.shadow.inherit = False
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGB(color)
    if line:
        shape.line.color.rgb = RGB(color)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, anchor=None):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    if anchor is not None:
        tf.vertical_anchor = anchor
    return tb, tf


def style_run(run, family=FONT_BODY, size=14, bold=False, color=NEUTRAL_950, italic=False):
    run.font.name = family
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = RGB(color)
    # set east-asian/cs font too so PowerPoint doesn't silently fall back
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            from pptx.oxml import parse_xml
            from pptx.oxml.ns import nsmap
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", family)


def add_text(slide, left, top, width, height, text, family=FONT_BODY, size=14, bold=False,
             color=NEUTRAL_950, align=PP_ALIGN.LEFT, anchor=None, line_spacing=1.0, italic=False):
    tb, tf = add_textbox(slide, left, top, width, height, anchor=anchor)
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    style_run(run, family=family, size=size, bold=bold, color=color, italic=italic)
    return tb, tf


def add_logo(slide, variant_path, left, top, height):
    return slide.shapes.add_picture(str(variant_path), left, top, height=height)


def add_wordmark(slide, left, top, size=18, color=WHITE, dim_color=None):
    dim_color = dim_color or color
    tb, tf = add_textbox(slide, left, top, Cm(10), Cm(1.2))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r1 = p.add_run()
    r1.text = "alio"
    style_run(r1, family=FONT_DISPLAY, size=size, bold=True, color=color)
    r2 = p.add_run()
    r2.text = " analytics"
    style_run(r2, family=FONT_DISPLAY, size=size, bold=False, color=dim_color)
    return tb


def footer(slide, page_no=None, dark=False):
    color = NEUTRAL_500 if not dark else "9AA2B3"
    add_text(
        slide, MARGIN, SLIDE_H - Cm(1.1), Cm(10), Cm(0.8),
        "www.alio.ao", family=FONT_MONO, size=10, color=color,
    )
    if page_no is not None:
        add_text(
            slide, SLIDE_W - Cm(3) - MARGIN, SLIDE_H - Cm(1.1), Cm(3), Cm(0.8),
            str(page_no), family=FONT_MONO, size=10, color=color, align=PP_ALIGN.RIGHT,
        )
