"""
Generate Presentation.pptx (+ Presentation.potx) — Alio Analytics.
16:9. Master-style layouts (each built as a slide, since python-pptx cannot
create new custom Slide Layouts, only reuse/populate existing ones): Title,
Section divider, Content (bullets), Two-column, Data/chart, Quote, Team,
Closing/CTA. Logo + brand colors/fonts on every slide; tagline on title.
"""

import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent))

from pptx import Presentation
from pptx.util import Cm, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

from brand import (
    FONT_DISPLAY, FONT_BODY, FONT_MONO,
    PRIMARY_700, PRIMARY_950, ACCENT_500, ACCENT_600,
    NEUTRAL_950, NEUTRAL_700, NEUTRAL_500, NEUTRAL_200, NEUTRAL_50, WHITE, SUCCESS_500,
    TAGLINE_EN, TAGLINE_SECONDARY, COMPANY_NAME,
    LOGO_ICON_WHITE, LOGO_ICON_NAVY, OUT_DIR, hx,
)
from pptx_helpers import (
    SLIDE_W, SLIDE_H, MARGIN, add_slide, set_bg, add_gradient_rect, add_rect,
    add_text, add_textbox, add_logo, add_wordmark, footer, style_run, RGB,
)
from to_template import pptx_to_potx


def slide_title(prs):
    """Title layout — cover slide with tagline."""
    s = add_slide(prs)
    set_bg(s, PRIMARY_950)
    add_gradient_rect(s, 0, SLIDE_H - Cm(0.25), SLIDE_W, Cm(0.25))
    add_logo(s, LOGO_ICON_WHITE, MARGIN, Cm(1.3), height=Cm(1.6))
    add_wordmark(s, MARGIN + Cm(2.1), Cm(1.55), size=20, color=WHITE, dim_color="C4CAD6")

    add_text(s, MARGIN, Cm(7.2), Cm(24), Cm(3.2), "{{Título da Apresentação}}",
             family=FONT_DISPLAY, size=44, bold=True, color=WHITE)
    add_text(s, MARGIN, Cm(10.0), Cm(24), Cm(1.2), "{{Subtítulo — cliente / projeto / data}}",
             family=FONT_BODY, size=18, color="C4CAD6")
    add_text(s, MARGIN, SLIDE_H - Cm(2.0), Cm(20), Cm(0.9),
             TAGLINE_EN.upper().replace(".", ""), family=FONT_MONO, size=11, color=ACCENT_500)
    return s


def slide_section(prs, number="01", title_text="{{Título da Secção}}"):
    """Section divider layout — solid navy with a brand-gradient accent bar
    (python-pptx has no alpha/opacity control, so the SVG version's
    translucent navy-over-gradient wash is approximated with solid navy +
    a gradient edge rather than attempting an unreliable alpha blend)."""
    s = add_slide(prs)
    set_bg(s, PRIMARY_950)
    add_gradient_rect(s, SLIDE_W - Cm(0.3), 0, Cm(0.3), SLIDE_H)
    add_gradient_rect(s, 0, SLIDE_H - Cm(0.25), SLIDE_W, Cm(0.25))
    add_logo(s, LOGO_ICON_WHITE, MARGIN, Cm(1.2), height=Cm(1.3))
    add_text(s, MARGIN, Cm(6.0), Cm(6), Cm(4), number,
             family=FONT_DISPLAY, size=90, bold=True, color=ACCENT_500)
    add_text(s, MARGIN, Cm(10.2), Cm(26), Cm(2.5), title_text,
             family=FONT_DISPLAY, size=40, bold=True, color=WHITE)
    return s


def slide_content(prs):
    """Content layout — title + bullet list."""
    s = add_slide(prs)
    set_bg(s, WHITE)
    add_logo(s, LOGO_ICON_NAVY, MARGIN, Cm(0.9), height=Cm(1.0))
    add_text(s, MARGIN + Cm(1.4), Cm(1.0), Cm(10), Cm(0.9), "alio analytics",
              family=FONT_DISPLAY, size=13, bold=True, color=NEUTRAL_700)
    add_text(s, MARGIN, Cm(2.6), Cm(24), Cm(1.6), "{{Título do Slide}}",
             family=FONT_DISPLAY, size=32, bold=True, color=NEUTRAL_950)
    add_rect(s, MARGIN, Cm(4.15), Cm(3), Cm(0.12), ACCENT_500)

    tb, tf = add_textbox(s, MARGIN, Cm(5.2), Cm(28), Cm(11))
    bullets = [
        "{{Ponto de conteúdo um}}",
        "{{Ponto de conteúdo dois}}",
        "{{Ponto de conteúdo três}}",
        "{{Ponto de conteúdo quatro}}",
    ]
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(16)
        run = p.add_run()
        run.text = "●  " + b
        style_run(run, family=FONT_BODY, size=18, color=NEUTRAL_950)
    footer(s)
    return s


def slide_two_column(prs):
    """Two-column layout — title + two content panels."""
    s = add_slide(prs)
    set_bg(s, WHITE)
    add_logo(s, LOGO_ICON_NAVY, MARGIN, Cm(0.9), height=Cm(1.0))
    add_text(s, MARGIN, Cm(2.6), Cm(24), Cm(1.6), "{{Título do Slide}}",
             family=FONT_DISPLAY, size=32, bold=True, color=NEUTRAL_950)
    add_rect(s, MARGIN, Cm(4.15), Cm(3), Cm(0.12), ACCENT_500)

    col_w = Cm(13.6)
    col_y = Cm(5.4)
    col_h = Cm(11.5)
    for i, (heading, color) in enumerate((("{{Coluna A}}", PRIMARY_700), ("{{Coluna B}}", ACCENT_600))):
        x = MARGIN + i * (col_w + Cm(1))
        add_rect(s, x, col_y, col_w, Cm(0.1), color)
        add_text(s, x, col_y + Cm(0.3), col_w, Cm(1), heading,
                 family=FONT_DISPLAY, size=20, bold=True, color=NEUTRAL_950)
        add_text(s, x, col_y + Cm(1.4), col_w, col_h - Cm(1.4),
                 "{{Texto ou pontos desta coluna}}", family=FONT_BODY, size=15, color=NEUTRAL_700)
    footer(s)
    return s


def slide_data(prs):
    """Data/chart layout — title + a real embedded chart on brand colors."""
    s = add_slide(prs)
    set_bg(s, WHITE)
    add_logo(s, LOGO_ICON_NAVY, MARGIN, Cm(0.9), height=Cm(1.0))
    add_text(s, MARGIN, Cm(2.6), Cm(24), Cm(1.6), "{{Título dos Dados}}",
             family=FONT_DISPLAY, size=32, bold=True, color=NEUTRAL_950)
    add_rect(s, MARGIN, Cm(4.15), Cm(3), Cm(0.12), ACCENT_500)

    chart_data = CategoryChartData()
    chart_data.categories = ["{{Cat 1}}", "{{Cat 2}}", "{{Cat 3}}", "{{Cat 4}}"]
    chart_data.add_series("{{Série}}", (4.2, 6.8, 5.1, 8.4))
    x, y, cx, cy = MARGIN, Cm(5.4), Cm(28), Cm(11.5)
    gframe = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data)
    chart = gframe.chart
    chart.has_legend = False
    plot = chart.plots[0]
    plot.has_data_labels = False
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = RGB(PRIMARY_700)
    try:
        chart.font.name = FONT_BODY
        chart.font.size = Pt(12)
        chart.font.color.rgb = RGB(NEUTRAL_700)
    except Exception:
        pass
    footer(s)
    return s


def slide_quote(prs):
    """Quote layout — large pull-quote on brand gradient."""
    s = add_slide(prs)
    add_gradient_rect(s, 0, 0, SLIDE_W, SLIDE_H, angle=90)
    add_text(s, Cm(3), Cm(3.5), Cm(2.5), Cm(2.5), "“",
             family=FONT_DISPLAY, size=90, bold=True, color=WHITE)
    add_text(s, Cm(3), Cm(6.8), Cm(28), Cm(6), "{{Citação de destaque do cliente ou parceiro.}}",
             family=FONT_DISPLAY, size=30, bold=True, color=WHITE, line_spacing=1.15)
    add_text(s, Cm(3), Cm(13.6), Cm(20), Cm(1), "{{Nome}} · {{Cargo, Empresa}}",
             family=FONT_BODY, size=16, color="C4CAD6")
    return s


def slide_team(prs):
    """Team layout — logo + a row of placeholder member cards."""
    s = add_slide(prs)
    set_bg(s, "F7F8FA")
    add_logo(s, LOGO_ICON_NAVY, MARGIN, Cm(0.9), height=Cm(1.0))
    add_text(s, MARGIN, Cm(2.6), Cm(24), Cm(1.6), "{{Equipa do Projeto}}",
             family=FONT_DISPLAY, size=32, bold=True, color=NEUTRAL_950)
    add_rect(s, MARGIN, Cm(4.15), Cm(3), Cm(0.12), ACCENT_500)

    card_w, card_h, gap = Cm(6.6), Cm(9.5), Cm(0.8)
    for i in range(4):
        x = MARGIN + i * (card_w + gap)
        y = Cm(5.4)
        add_rect(s, x, y, card_w, card_h, WHITE)
        add_rect(s, x, y, card_w, Cm(0.08), ACCENT_500)
        add_rect(s, x + Cm(2.05), y + Cm(1.0), Cm(2.5), Cm(2.5), "EEF0FB")
        add_text(s, x + Cm(0.4), y + Cm(3.9), card_w - Cm(0.8), Cm(0.9), "{{Nome}}",
                 family=FONT_DISPLAY, size=15, bold=True, color=NEUTRAL_950, align=PP_ALIGN.CENTER)
        add_text(s, x + Cm(0.4), y + Cm(4.6), card_w - Cm(0.8), Cm(0.9), "{{Cargo}}",
                 family=FONT_BODY, size=12, color=NEUTRAL_500, align=PP_ALIGN.CENTER)
    footer(s)
    return s


def slide_closing(prs):
    """Closing / CTA layout."""
    s = add_slide(prs)
    set_bg(s, PRIMARY_950)
    add_gradient_rect(s, 0, 0, SLIDE_W, Cm(0.25))
    add_logo(s, LOGO_ICON_WHITE, MARGIN, Cm(7.6), height=Cm(1.8))
    add_wordmark(s, MARGIN, Cm(9.5), size=22, color=WHITE, dim_color="C4CAD6")
    add_text(s, MARGIN, Cm(10.8), Cm(24), Cm(1), "{{Próximos passos}} — {{contacto@alio.ao}}",
             family=FONT_BODY, size=16, color="C4CAD6")
    add_text(s, MARGIN, SLIDE_H - Cm(2.0), Cm(20), Cm(0.9),
             TAGLINE_SECONDARY.upper().replace(".", ""), family=FONT_MONO, size=10, color=ACCENT_500)
    return s


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs)
    slide_section(prs, "01", "{{Título da Secção}}")
    slide_content(prs)
    slide_two_column(prs)
    slide_data(prs)
    slide_quote(prs)
    slide_team(prs)
    slide_closing(prs)

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    pptx_path = OUT_DIR / "alio-presentation.pptx"
    potx_path = OUT_DIR / "alio-presentation.potx"
    prs.save(pptx_path)
    pptx_to_potx(pptx_path, potx_path)
    print("  OK", pptx_path.name, "+", potx_path.name, "-", len(prs.slides), "slides")


if __name__ == "__main__":
    build()
